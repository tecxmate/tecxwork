#!/usr/bin/env python3
"""Generate the PowerPoint deck from the manual's own source.

    python3 docs/manual/src/build_pptx.py            # English
    python3 docs/manual/src/build_pptx.py zh vi      # plus 繁中 and Tiếng Việt

Content is parsed out of manual.src.html rather than retyped, so the deck cannot drift from
the manual: every screen slide takes its step number, title, route, purpose and control list
straight from the same markup the HTML renders. Translated decks reuse the same i18n
dictionaries, keyed by the data-t attributes.

Two things python-pptx cannot do, handled here:
  - it has no WebP decoder, so screenshots are converted to JPEG on the way in;
  - it cannot embed SVG, so the two explanatory diagrams are pre-rendered to PNG by
    render-diagrams.mjs and picked up from /tmp/deck-assets.
"""
from __future__ import annotations

import io
import json
import os
import re
import sys

from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "manual.src.html")
SHOTS = os.path.join(HERE, "screenshots")
STRINGS = os.path.join(HERE, "strings")
DIAGRAMS = "/tmp/deck-assets"
OUTDIR = os.path.normpath(os.path.join(HERE, "..", "..", "..", "public"))

# 16:9
W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x1D, 0x1D, 0x1F)
INK2 = RGBColor(0x4A, 0x47, 0x50)
MUTED = RGBColor(0x6E, 0x6E, 0x73)
GROUND = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE5, 0xE1, 0xEC)
ACCENT = RGBColor(0x8C, 0x52, 0xFF)

ROLE = {  # matches the manual's role colours
    "1": RGBColor(0x0F, 0x76, 0x6E),   # applicant
    "2": RGBColor(0x24, 0x5E, 0x9D),   # employer
    "3": RGBColor(0x8C, 0x52, 0xFF),   # agency
    "4": RGBColor(0x9A, 0x5B, 0x12),   # admin
}

# PowerPoint substitutes fonts silently; name a CJK-capable face explicitly so the
# Chinese deck does not fall back to something with no glyphs.
FONT_LATIN = "Helvetica Neue"
FONT_EA = {"zh": "PingFang TC", "vi": "Helvetica Neue", "en": "Helvetica Neue"}
FONT_DISPLAY = {"zh": "PingFang TC", "vi": "Georgia", "en": "Georgia"}


# ----------------------------------------------------------------- content ----
def text_of(node, strings, lang):
    """Visible text for a node, swapped to `lang` when it carries a data-t key."""
    if node is None:
        return ""
    key = node.get("data-t") if hasattr(node, "get") else None
    if key and lang != "en" and key in strings:
        html = strings[key]
        return re.sub(r"\s+", " ", BeautifulSoup(html, "html.parser").get_text(" ")).strip()
    return re.sub(r"\s+", " ", node.get_text(" ")).strip()


def route_of(sec) -> str:
    """The route line, with any trailing pill separated rather than run together."""
    el = sec.select_one(".route")
    if el is None:
        return ""
    pill = el.select_one(".pill")
    note = ""
    if pill is not None:
        note = re.sub(r"\s+", " ", pill.get_text(" ")).strip()
        pill.extract()
    base = re.sub(r"\s+", " ", el.get_text(" ")).strip()
    return base + ("   ·   " + note if note else "")


def parse(lang: str, strings: dict) -> dict:
    with open(SRC, encoding="utf-8") as fh:
        soup = BeautifulSoup(fh.read(), "html.parser")

    def t(node):
        return text_of(node, strings, lang)

    acts = []
    for div in soup.select('div.act[id^="act"]'):
        acts.append({
            "id": div["id"],
            "no": t(div.select_one(".act-no")),
            "title": t(div.select_one("h2")),
            "blurb": t(div.select_one("p:not(.act-no)")),
            "who": t(div.select_one(".who")),
            "role": div["id"][-1],
        })

    screens = []
    for sec in soup.select("section.screen"):
        step_el = sec.select_one(".step")
        if not step_el:
            continue
        step = step_el.get_text(strip=True)
        img = sec.select_one("figure img")
        m = re.search(r"\{\{IMG:([a-z0-9\-]+)\}\}", str(sec))
        controls = []
        for dl in sec.select("dl.controls"):
            for dt in dl.select("dt"):
                dd = dt.find_next_sibling("dd")
                if dd is not None:
                    controls.append((t(dt), t(dd)))
        screens.append({
            "step": step,
            "act": step.split(".")[0] if "." in step else "0",
            "title": t(sec.select_one(".screen-head h3")),
            "route": route_of(sec),
            "purpose": t(sec.select_one(".purpose")),
            "shot": m.group(1) if m else None,
            "controls": controls,
        })

    hero = soup.select_one(".hero")
    return {
        "title": t(soup.select_one(".hero h1")),
        "lede": t(hero.select_one(".lede")) if hero else "",
        "eyebrow": t(hero.select_one(".eyebrow")) if hero else "",
        "acts": acts,
        "screens": screens,
        "trio": [(t(d.select_one("h5")), t(d.select_one("p"))) for d in soup.select(".trio > div")],
        "lifecycle": t(soup.select_one("#lifecycle h2")),
        "lifecycle_blurb": t(soup.select_one("#lifecycle .act p:not(.act-no)")),
        "jobfair": t(soup.select_one("#jobfair h2")),
        "jobfair_blurb": t(soup.select_one("#jobfair .act p:not(.act-no)")),
        "what": t(soup.select_one("#what h2")),
    }


# ------------------------------------------------------------------ layout ----
def jpeg(name: str, max_aspect: float = 1.05):
    """WebP -> JPEG bytes. Very tall pages are cropped so the slide image stays legible."""
    im = Image.open(os.path.join(SHOTS, name + ".webp")).convert("RGB")
    w, h = im.size
    cropped = False
    if h / w > max_aspect:
        im = im.crop((0, 0, w, int(w * max_aspect)))
        cropped = True
    if im.width > 1800:
        im = im.resize((1800, round(im.height * 1800 / im.width)), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, "JPEG", quality=82, optimize=True)
    buf.seek(0)
    return buf, cropped


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, *, bold=False, font=None, ea=None, space_after=6,
         first=False, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font or FONT_LATIN
    if ea:  # east-asian face, so CJK glyphs resolve
        rpr = run._r.get_or_add_rPr()
        from pptx.oxml.ns import qn
        ea_el = rpr.makeelement(qn("a:ea"), {"typeface": ea})
        rpr.append(ea_el)
    return p


def bg(slide, color=GROUND):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rule(slide, x, y, w, color=LINE, h=Emu(9525)):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def badge(slide, x, y, text, color):
    from pptx.enum.shapes import MSO_SHAPE
    w = Inches(0.42 + 0.11 * max(0, len(text) - 3))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    s.adjustments[0] = 0.18
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = "Menlo"
    return s


def fit(img_w, img_h, box_w, box_h):
    scale = min(box_w / img_w, box_h / img_h)
    return int(img_w * scale), int(img_h * scale)


# ------------------------------------------------------------------- build ----
def build(lang: str) -> str:
    strings = {}
    if lang != "en":
        with open(os.path.join(STRINGS, lang + ".json"), encoding="utf-8") as fh:
            strings = json.load(fh)
    d = parse(lang, strings)
    ea = FONT_EA[lang]
    disp = FONT_DISPLAY[lang]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    # ---- title ----
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    rule(s, Inches(0.9), Inches(2.05), Inches(2.2), ACCENT, Emu(38100))
    tf = textbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.5))
    para(tf, d["title"], 46, INK, first=True, font=disp, ea=ea, space_after=2)
    tf2 = textbox(s, Inches(0.9), Inches(3.5), Inches(9.6), Inches(2.2))
    para(tf2, d["lede"], 15, INK2, first=True, ea=ea, space_after=10)
    tf3 = textbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5))
    para(tf3, d["eyebrow"], 11, MUTED, first=True, font="Menlo", ea=ea)

    # ---- what the platform is (three products) ----
    s = prs.slides.add_slide(blank); bg(s)
    tf = textbox(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.9))
    para(tf, d["what"], 30, INK, first=True, font=disp, ea=ea)
    rule(s, Inches(0.9), Inches(1.55), Inches(11.5))
    colw = Inches(3.6)
    for i, (h5, body) in enumerate(d["trio"][:3]):
        x = Inches(0.9 + i * 3.87)
        rule(s, x, Inches(1.95), colw, list(ROLE.values())[i], Emu(28575))
        tfc = textbox(s, x, Inches(2.15), colw, Inches(4.4))
        para(tfc, h5, 17, list(ROLE.values())[i], bold=True, first=True, ea=ea, space_after=8)
        para(tfc, body, 12.5, INK2, ea=ea)

    # ---- the two diagrams ----
    for png, heading, blurb in (
        ("diagram-lifecycle.png", d["lifecycle"], d["lifecycle_blurb"]),
        ("diagram-jobfair.png", d["jobfair"], d["jobfair_blurb"]),
    ):
        path = os.path.join(DIAGRAMS, png)
        if not os.path.exists(path):
            print("  ! missing diagram %s — run render-diagrams.mjs" % png)
            continue
        s = prs.slides.add_slide(blank); bg(s, WHITE)
        tf = textbox(s, Inches(0.75), Inches(0.34), Inches(11.8), Inches(0.9))
        para(tf, heading, 25, INK, first=True, font=disp, ea=ea, space_after=4)
        para(tf, blurb, 11, MUTED, ea=ea)
        iw, ih = Image.open(path).size
        nw, nh = fit(iw, ih, Inches(11.9), Inches(5.55))
        s.shapes.add_picture(path, int((W - nw) / 2), Inches(1.55), nw, nh)

    # ---- acts and their screens ----
    for act in d["acts"]:
        colour = ROLE[act["role"]]
        s = prs.slides.add_slide(blank); bg(s, WHITE)
        rule(s, Inches(0.9), Inches(2.4), Inches(1.6), colour, Emu(38100))
        tf = textbox(s, Inches(0.9), Inches(2.65), Inches(11.0), Inches(2.2))
        para(tf, act["no"].upper(), 13, colour, bold=True, first=True, font="Menlo",
             ea=ea, space_after=10)
        para(tf, act["title"], 34, INK, font=disp, ea=ea, space_after=12)
        para(tf, act["blurb"], 13, INK2, ea=ea)

        for sc in [x for x in d["screens"] if x["act"] == act["role"]]:
            s = prs.slides.add_slide(blank); bg(s)

            # Screenshot right, full-bleed to the slide edges; text left. The screenshot is
            # the substance of the slide, so it gets the larger half and no margin — in the
            # earlier layout it sat in a 7.15in box under a header and was too small to read.
            PANEL_X, PANEL_W = Inches(5.05), Inches(8.283)
            if sc["shot"]:
                buf, cropped = jpeg(sc["shot"])
                iw, ih = Image.open(buf).size
                buf.seek(0)
                nw, nh = fit(iw, ih, PANEL_W, H)
                s.shapes.add_picture(buf, PANEL_X + int((PANEL_W - nw) / 2),
                                     int((H - nh) / 2), nw, nh)

            TX, TW = Inches(0.62), Inches(4.05)
            badge(s, TX, Inches(0.62), sc["step"], colour)
            tf = textbox(s, TX, Inches(1.14), TW, Inches(1.0))
            para(tf, sc["title"], 20, INK, first=True, font=disp, ea=ea, space_after=3)
            if sc["route"]:
                para(tf, sc["route"], 9.5, MUTED, font="Menlo", ea=ea)

            tfp = textbox(s, TX, Inches(2.35), TW, Inches(4.75))
            para(tfp, sc["purpose"], 11, INK2, first=True, ea=ea, space_after=11)
            for label, desc in sc["controls"][:5]:
                p = tfp.add_paragraph()
                p.space_after = Pt(6)
                r1 = p.add_run(); r1.text = label + " — "
                r1.font.size = Pt(9.5); r1.font.bold = True; r1.font.color.rgb = INK
                r1.font.name = FONT_LATIN
                r2 = p.add_run(); r2.text = desc if len(desc) < 135 else desc[:132] + "…"
                r2.font.size = Pt(9.5); r2.font.color.rgb = INK2; r2.font.name = FONT_LATIN
                if ea:
                    from pptx.oxml.ns import qn
                    for r in (r1, r2):
                        rpr = r._r.get_or_add_rPr()
                        rpr.append(rpr.makeelement(qn("a:ea"), {"typeface": ea}))
            if len(sc["controls"]) > 5:
                para(tfp, "+ %d more in the manual" % (len(sc["controls"]) - 5), 9, MUTED,
                     font="Menlo", ea=ea, italic=True)
            if sc["shot"] and cropped:
                # note it in the text column, not over the screenshot
                tfc = textbox(s, TX, Inches(6.85), TW, Inches(0.3))
                para(tfc, "screenshot shows the top of the page", 8.5, MUTED, first=True,
                     font="Menlo", ea=ea)

    # ---- closing ----
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    tf = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.0))
    para(tf, d["title"], 30, INK, first=True, font=disp, ea=ea, space_after=14)
    para(tf, "yangluck.tecxmate.com/documentation", 17, ACCENT, ea=ea, space_after=10)
    para(tf, "The full manual — every control on every screen, in English, 繁體中文 and "
             "Tiếng Việt — lives at the link above and works offline.", 12.5, MUTED, ea=ea)

    out = os.path.join(OUTDIR, "tecxwork-manual%s.pptx" % ("" if lang == "en" else "-" + lang))
    prs.save(out)
    mb = os.path.getsize(out) / 1024 / 1024
    print("%-34s %3d slides  %5.1f MB" % (os.path.basename(out), len(prs.slides.__iter__.__self__._sldIdLst), mb))
    return out


if __name__ == "__main__":
    langs = sys.argv[1:] or ["en"]
    for lang in langs:
        build(lang)
